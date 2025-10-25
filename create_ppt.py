from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "EADA Pro - Phase 1 Review"
subtitle.text = "Edge AI for Adaptive TV/Billboard Applications\nIntelligent Brightness, Volume & Media Control\nOctober 2025"

# Slide 2: Project Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Project Overview'
tf = body_shape.text_frame
tf.text = 'Vision: Intelligent, privacy-preserving adaptive system for TV/billboard applications'

p = tf.add_paragraph()
p.text = '• Real-time AI adjusts brightness, volume, and media based on audience presence'
p.level = 0

p = tf.add_paragraph()
p.text = '• Edge-first deployment with computer vision and audio processing'
p.level = 0

p = tf.add_paragraph()
p.text = '• Privacy-preserving: No data storage, local processing only'
p.level = 0

p = tf.add_paragraph()
p.text = '• Phase 1 Focus: Core functionality - distance detection, adaptive controls, media pause/resume'
p.level = 0

# Slide 3: Technical Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Technical Architecture'
tf = body_shape.text_frame
tf.text = 'Modular Design with Clean Separation of Concerns'

p = tf.add_paragraph()
p.text = 'Perception Layer:'
p.level = 0
p = tf.add_paragraph()
p.text = '• Face Detector (MediaPipe Face Mesh)'
p.level = 1
p = tf.add_paragraph()
p.text = '• Face Counter with position tracking'
p.level = 1

p = tf.add_paragraph()
p.text = 'Intelligence Layer:'
p.level = 0
p = tf.add_paragraph()
p.text = '• Environment Monitor (ambient light)'
p.level = 1
p = tf.add_paragraph()
p.text = '• Audio Analyzer (background music)'
p.level = 1

p = tf.add_paragraph()
p.text = 'Adaptation Layer:'
p.level = 0
p = tf.add_paragraph()
p.text = '• Brightness Controller (30-100% range)'
p.level = 1
p = tf.add_paragraph()
p.text = '• Volume Controller (20-100% with steep curve)'
p.level = 1
p = tf.add_paragraph()
p.text = '• Weighted Adapter (multi-face support)'
p.level = 1

p = tf.add_paragraph()
p.text = 'Core: System Manager orchestrates all modules'
p.level = 0

# Slide 4: Key Algorithms
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Key Algorithms Implemented'
tf = body_shape.text_frame
tf.text = 'Distance Estimation & Adaptive Control'

p = tf.add_paragraph()
p.text = 'Distance Calculation:'
p.level = 0
p = tf.add_paragraph()
p.text = '• Face width triangulation using MediaPipe landmarks (234-454)'
p.level = 1
p = tf.add_paragraph()
p.text = '• Calibrated for accurate real-world measurements'
p.level = 1

p = tf.add_paragraph()
p.text = 'Adaptive Logic:'
p.level = 0
p = tf.add_paragraph()
p.text = '• Inverse relationship: Farther = brighter/louder, Closer = dimmer/quieter'
p.level = 1
p = tf.add_paragraph()
p.text = '• Brightness: Linear mapping 30-100%'
p.level = 1
p = tf.add_paragraph()
p.text = '• Volume: Steep curve (power 0.4) for noticeable changes'
p.level = 1

p = tf.add_paragraph()
p.text = 'Stability Features:'
p.level = 0
p = tf.add_paragraph()
p.text = '• Waits 1.5s for distance to stabilize'
p.level = 1
p = tf.add_paragraph()
p.text = '• Movement threshold: 5cm resets timer'
p.level = 1
p = tf.add_paragraph()
p.text = '• Grace range: ±10cm prevents micro-adjustments'
p.level = 1

# Slide 5: Features Implemented
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Phase 1 Features Completed'
tf = body_shape.text_frame
tf.text = 'Core Functionality Delivered'

p = tf.add_paragraph()
p.text = '✅ Distance-based Brightness Control'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Distance-based Volume Control with Steep Curve'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Presence-based Media Pause/Resume (3s timeout)'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Face Detection with MediaPipe Face Mesh'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Stability Logic for Crowd Scenarios'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Grace Range to Prevent Constant Adjustments'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Clean Console Output (Suppressed Warnings)'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Python 3.10 Compatibility with MediaPipe'
p.level = 0

# Slide 6: Code Structure
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Code Structure & Organization'
tf = body_shape.text_frame
tf.text = 'Clean, Modular Python Implementation'

p = tf.add_paragraph()
p.text = 'src/config/settings.py - Centralized configuration'
p.level = 0
p = tf.add_paragraph()
p.text = 'src/modules/perception/ - Face detection & counting'
p.level = 0
p = tf.add_paragraph()
p.text = 'src/modules/adaptation/ - Brightness & volume controllers'
p.level = 0
p = tf.add_paragraph()
p.text = 'src/modules/intelligence/ - Environment & audio analysis'
p.level = 0
p = tf.add_paragraph()
p.text = 'src/core/system_manager.py - Main orchestration'
p.level = 0
p = tf.add_paragraph()
p.text = 'src/main.py - Clean startup with logging control'
p.level = 0

# Slide 7: Challenges Overcome
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Challenges Overcome'
tf = body_shape.text_frame
tf.text = 'Technical Hurdles & Solutions'

p = tf.add_paragraph()
p.text = 'Python Version Compatibility:'
p.level = 0
p = tf.add_paragraph()
p.text = '• MediaPipe 0.10.9 requires Python ≤3.11'
p.level = 1
p = tf.add_paragraph()
p.text = '• Solution: Created Python 3.10 virtual environment'
p.level = 1

p = tf.add_paragraph()
p.text = 'Stability in Crowd Scenarios:'
p.level = 0
p = tf.add_paragraph()
p.text = '• Constant adjustments with moving crowds'
p.level = 1
p = tf.add_paragraph()
p.text = '• Solution: Stability timer + grace range + movement threshold'
p.level = 1

p = tf.add_paragraph()
p.text = 'Console Output Management:'
p.level = 0
p = tf.add_paragraph()
p.text = '• MediaPipe INFO messages cluttering output'
p.level = 1
p = tf.add_paragraph()
p.text = '• Solution: Suppressed logging to WARNING+ level'
p.level = 1

p = tf.add_paragraph()
p.text = 'Volume Control Sensitivity:'
p.level = 0
p = tf.add_paragraph()
p.text = '• Linear volume changes barely noticeable'
p.level = 1
p = tf.add_paragraph()
p.text = '• Solution: Implemented steep power curve (0.4)'
p.level = 1

# Slide 8: Demo Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Demo & Validation Results'
tf = body_shape.text_frame
tf.text = 'System Performance & Behavior'

p = tf.add_paragraph()
p.text = '✅ System starts cleanly with clear status messages'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Face detection works reliably with MediaPipe Face Mesh'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Brightness adjusts smoothly based on distance'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Volume changes noticeably with steep curve'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Media pauses after 3s absence, resumes on return'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Stability logic prevents erratic adjustments in crowds'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Grace range eliminates micro-adjustments'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Console remains clean during operation'
p.level = 0

# Slide 9: Future Phases
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Future Development Roadmap'
tf = body_shape.text_frame
tf.text = 'Phase 2-6 Expansion Plans'

p = tf.add_paragraph()
p.text = 'Phase 2: Gesture Integration'
p.level = 0
p = tf.add_paragraph()
p.text = '• Hand detection and gesture-to-action mapping'
p.level = 1
p = tf.add_paragraph()
p.text = '• Thumb-index distance for volume, wrist position for brightness'
p.level = 1

p = tf.add_paragraph()
p.text = 'Phase 3: API + Dashboard'
p.level = 0
p = tf.add_paragraph()
p.text = '• FastAPI backend with live metrics'
p.level = 1
p = tf.add_paragraph()
p.text = '• React dashboard with face counting display'
p.level = 1

p = tf.add_paragraph()
p.text = 'Phase 4-6: Enterprise & Public Display Features'
p.level = 0
p = tf.add_paragraph()
p.text = '• Crowd analysis, weather adaptation, energy optimization'
p.level = 1
p = tf.add_paragraph()
p.text = '• Security monitoring, analytics, edge deployment'
p.level = 1

# Slide 10: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Conclusion & Achievements'
tf = body_shape.text_frame
tf.text = 'Phase 1 Success & Next Steps'

p = tf.add_paragraph()
p.text = '🎯 Phase 1 Objectives Met:'
p.level = 0
p = tf.add_paragraph()
p.text = '• Core adaptive functionality implemented and tested'
p.level = 1
p = tf.add_paragraph()
p.text = '• Robust distance-based controls with stability features'
p.level = 1
p = tf.add_paragraph()
p.text = '• Clean, maintainable codebase ready for expansion'
p.level = 1

p = tf.add_paragraph()
p.text = '🔧 Technical Excellence:'
p.level = 0
p = tf.add_paragraph()
p.text = '• Proper virtual environment management'
p.level = 1
p = tf.add_paragraph()
p.text = '• Modular architecture for easy feature addition'
p.level = 1
p = tf.add_paragraph()
p.text = '• Comprehensive error handling and logging'
p.level = 1

p = tf.add_paragraph()
p.text = '🚀 Ready for Phase 2: Gesture recognition integration'
p.level = 0

# Save the presentation
prs.save('EADA_Pro_Phase1_Review.pptx')
print("PPT created successfully: EADA_Pro_Phase1_Review.pptx")